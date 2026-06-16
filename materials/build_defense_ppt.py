from __future__ import annotations

from pathlib import Path

from pptx import Presentation
from pptx.dml.color import RGBColor
from pptx.enum.shapes import MSO_AUTO_SHAPE_TYPE
from pptx.enum.text import MSO_ANCHOR, PP_ALIGN
from pptx.oxml.ns import qn
from pptx.util import Cm, Pt


BASE_DIR = Path(__file__).resolve().parent
TEMPLATE = BASE_DIR / "template_source.pptx"
OUTPUT = BASE_DIR / "hotel_defense_presentation_liu_yang.pptx"
TEMPLATE_MEDIA = BASE_DIR / "template_unzip" / "ppt" / "media"
THESIS_MEDIA = BASE_DIR / "thesis_media"

FONT_CN = "Microsoft YaHei"
FONT_EN = "Aptos"

BLUE = RGBColor(22, 88, 160)
BLUE_DARK = RGBColor(13, 63, 120)
RED = RGBColor(194, 37, 37)
GOLD = RGBColor(201, 156, 61)
TEXT = RGBColor(36, 42, 48)
MUTED = RGBColor(112, 118, 126)
WHITE = RGBColor(255, 255, 255)
PANEL = RGBColor(255, 255, 255)
LINE = RGBColor(215, 223, 233)
SOFT_BLUE = RGBColor(238, 245, 252)


def set_font(run, size: float, bold: bool = False, color: RGBColor = TEXT) -> None:
    run.font.size = Pt(size)
    run.font.bold = bold
    run.font.name = FONT_EN
    run.font.color.rgb = color
    rpr = run._r.get_or_add_rPr()
    rpr.set(qn("a:ea"), FONT_CN)
    rpr.set(qn("a:latin"), FONT_EN)


def clear_tf(tf) -> None:
    tf.clear()
    tf.word_wrap = True
    tf.vertical_anchor = MSO_ANCHOR.TOP


def add_para(tf, text: str, *, size: float = 12, bold: bool = False, color: RGBColor = TEXT,
             bullet: bool = False, align=PP_ALIGN.LEFT, after: float = 4) -> None:
    p = tf.paragraphs[0] if len(tf.paragraphs) == 1 and not tf.paragraphs[0].text else tf.add_paragraph()
    p.text = text
    p.alignment = align
    p.space_after = Pt(after)
    p.line_spacing = 1.12
    if bullet:
        p.bullet = True
    for run in p.runs:
        set_font(run, size, bold=bold, color=color)


def add_bg(slide, bg_name: str) -> None:
    slide.shapes.add_picture(str(TEMPLATE_MEDIA / bg_name), 0, 0, width=slide.part.slide_layout.part.package.presentation_part.presentation.slide_width, height=slide.part.slide_layout.part.package.presentation_part.presentation.slide_height)


def add_textbox(slide, left, top, width, height):
    return slide.shapes.add_textbox(left, top, width, height)


def add_panel(slide, left, top, width, height, *, fill=PANEL, line=LINE, rounded=True, transparency: float = 0.08):
    shp_type = MSO_AUTO_SHAPE_TYPE.ROUNDED_RECTANGLE if rounded else MSO_AUTO_SHAPE_TYPE.RECTANGLE
    shp = slide.shapes.add_shape(shp_type, left, top, width, height)
    shp.fill.solid()
    shp.fill.fore_color.rgb = fill
    shp.fill.transparency = transparency
    shp.line.color.rgb = line
    shp.line.width = Pt(0.8)
    return shp


def add_title_plain(slide, title: str, subtitle: str | None = None) -> None:
    tb = add_textbox(slide, Cm(1.0), Cm(1.05), Cm(18.0), Cm(0.8))
    clear_tf(tb.text_frame)
    add_para(tb.text_frame, title, size=21, bold=True, color=BLUE_DARK, after=0)
    if subtitle:
        sb = add_textbox(slide, Cm(1.02), Cm(1.72), Cm(18.0), Cm(0.42))
        clear_tf(sb.text_frame)
        add_para(sb.text_frame, subtitle, size=9.3, color=MUTED, after=0)


def add_section_title(slide, title: str, subtitle: str) -> None:
    tb = add_textbox(slide, Cm(3.0), Cm(4.2), Cm(12.5), Cm(2.2))
    clear_tf(tb.text_frame)
    add_para(tb.text_frame, title, size=24, bold=True, color=WHITE, align=PP_ALIGN.CENTER, after=8)
    add_para(tb.text_frame, subtitle, size=11.5, color=WHITE, align=PP_ALIGN.CENTER, after=0)


def add_pic(slide, name: str, left, top, width=None, height=None):
    slide.shapes.add_picture(str(THESIS_MEDIA / name), left, top, width=width, height=height)


def delete_all_slides(prs: Presentation) -> None:
    for idx in range(len(prs.slides) - 1, -1, -1):
        slide_id = prs.slides._sldIdLst[idx]
        prs.part.drop_rel(slide_id.rId)
        del prs.slides._sldIdLst[idx]


def validate(prs: Presentation) -> None:
    for slide in prs.slides:
        for shape in slide.shapes:
            if shape.left + shape.width > prs.slide_width + 1:
                raise ValueError("shape exceeds width")
            if shape.top + shape.height > prs.slide_height + 1:
                raise ValueError("shape exceeds height")


def build_cover(prs: Presentation) -> None:
    slide = prs.slides.add_slide(prs.slide_layouts[0])
    add_bg(slide, "image1.png")

    title = add_textbox(slide, Cm(1.45), Cm(1.55), Cm(21.0), Cm(3.6))
    clear_tf(title.text_frame)
    add_para(title.text_frame, "基于SpringBoot与RAG的", size=23, bold=True, color=WHITE, after=2)
    add_para(title.text_frame, "酒店运营支持系统设计与实现", size=24, bold=True, color=WHITE, after=0)

    sub = add_textbox(slide, Cm(1.5), Cm(5.0), Cm(8.5), Cm(0.6))
    clear_tf(sub.text_frame)
    add_para(sub.text_frame, "本科毕业论文答辩", size=13, bold=True, color=BLUE_DARK, after=0)

    panel = add_panel(slide, Cm(1.35), Cm(9.0), Cm(9.2), Cm(3.4), fill=WHITE, transparency=0.02)
    tf = panel.text_frame
    clear_tf(tf)
    add_para(tf, "答辩信息", size=14.5, bold=True, color=BLUE_DARK, after=6)
    for item in [
        "学生：刘洋",
        "学号：2022105400197",
        "指导教师：刘会芝",
        "专业班级：人工智能 3 班",
    ]:
        add_para(tf, item, size=11.4, color=TEXT, after=3)

    panel2 = add_panel(slide, Cm(11.0), Cm(9.0), Cm(10.7), Cm(3.4), fill=SOFT_BLUE, transparency=0.03)
    tf2 = panel2.text_frame
    clear_tf(tf2)
    add_para(tf2, "答辩关键词", size=14.5, bold=True, color=BLUE_DARK, after=6)
    for item in [
        "运营看板 + 知识管理 + RAG问答",
        "入住率 / 营收 / 订单需求分类预测",
        "决策支持与测试验证闭环",
    ]:
        add_para(tf2, item, size=11.2, bullet=True, color=TEXT, after=3)


def build_toc(prs: Presentation) -> None:
    slide = prs.slides.add_slide(prs.slide_layouts[0])
    add_bg(slide, "image2.png")

    panel = add_panel(slide, Cm(10.4), Cm(1.5), Cm(12.8), Cm(10.5), fill=WHITE, transparency=0.18)
    tf = panel.text_frame
    clear_tf(tf)
    add_para(tf, "答辩目录", size=18, bold=True, color=BLUE_DARK, after=10)
    items = [
        "01  研究背景与本文切入",
        "02  系统目标与总体设计",
        "03  系统实现与页面展示",
        "04  实验结果与测试验证",
        "05  总结与展望",
    ]
    for item in items:
        add_para(tf, item, size=13.4, color=TEXT, after=8)


def build_section(prs: Presentation, title: str, subtitle: str) -> None:
    slide = prs.slides.add_slide(prs.slide_layouts[0])
    add_bg(slide, "image3.png")
    add_section_title(slide, title, subtitle)


def build_background_gap(prs: Presentation) -> None:
    slide = prs.slides.add_slide(prs.slide_layouts[0])
    add_bg(slide, "image4.png")
    add_title_plain(slide, "研究背景与本文切入", "强调现有研究缺口与本文补足方式")

    left = add_panel(slide, Cm(1.0), Cm(2.4), Cm(10.8), Cm(8.5), fill=WHITE, transparency=0.02)
    tf = left.text_frame
    clear_tf(tf)
    add_para(tf, "现有研究不足", size=14.5, bold=True, color=BLUE_DARK, after=8)
    for item in [
        "多聚焦单一能力，如管理平台或知识问答。",
        "较少把运营分析、RAG 问答、预测支持整合为一体化平台。",
        "面向中小型酒店的轻量化运营支持方案仍较少。",
    ]:
        add_para(tf, item, size=11.6, bullet=True, color=TEXT, after=4)

    right = add_panel(slide, Cm(12.3), Cm(2.4), Cm(11.1), Cm(8.5), fill=SOFT_BLUE, transparency=0.04)
    tf2 = right.text_frame
    clear_tf(tf2)
    add_para(tf2, "本文补足方式", size=14.5, bold=True, color=BLUE_DARK, after=8)
    for item in [
        "采用 SpringBoot + Flask + Vue3 搭建完整工程链路。",
        "整合知识管理、RAG 问答、运营总览、三类预测和决策支持。",
        "通过测试与实验结果证明系统具备可运行、可展示、可验证价值。",
    ]:
        add_para(tf2, item, size=11.6, bullet=True, color=TEXT, after=4)


def build_goal_arch(prs: Presentation) -> None:
    slide = prs.slides.add_slide(prs.slide_layouts[0])
    add_bg(slide, "image4.png")
    add_title_plain(slide, "系统目标与总体架构", "回答“这套系统做了什么、怎么组织起来的”")

    goal = add_panel(slide, Cm(1.0), Cm(2.35), Cm(9.0), Cm(3.5), fill=WHITE, transparency=0.02)
    tf = goal.text_frame
    clear_tf(tf)
    add_para(tf, "系统目标", size=14.2, bold=True, color=BLUE_DARK, after=8)
    for item in [
        "统一运营入口：集中展示核心经营指标。",
        "提供知识服务：支持知识维护与可追溯问答。",
        "形成预测支持：输出风险提示与建议动作。",
    ]:
        add_para(tf, item, size=11.4, bullet=True, color=TEXT, after=4)

    add_pic(slide, "image20.png", Cm(10.6), Cm(2.35), width=Cm(12.3))
    cap = add_textbox(slide, Cm(10.8), Cm(7.25), Cm(11.8), Cm(0.45))
    clear_tf(cap.text_frame)
    add_para(cap.text_frame, "总体架构图", size=9.4, color=BLUE_DARK, align=PP_ALIGN.CENTER, after=0)

    note = add_panel(slide, Cm(1.0), Cm(8.1), Cm(21.9), Cm(2.7), fill=SOFT_BLUE, transparency=0.05)
    clear_tf(note.text_frame)
    add_para(note.text_frame, "技术路线", size=13.6, bold=True, color=BLUE_DARK, after=6)
    add_para(note.text_frame, "Vue3 前端 + SpringBoot 业务后端 + Flask 智能服务 + MySQL 数据层 + text2vec / DeepSeek 能力支撑。", size=11.2, color=TEXT, after=0)


def build_module_data(prs: Presentation) -> None:
    slide = prs.slides.add_slide(prs.slide_layouts[0])
    add_bg(slide, "image4.png")
    add_title_plain(slide, "模块设计与关键数据", "只保留答辩真正需要解释的模块、表与接口")

    add_pic(slide, "image4.png", Cm(1.0), Cm(2.4), width=Cm(8.4))
    cap = add_textbox(slide, Cm(1.0), Cm(6.7), Cm(8.4), Cm(0.45))
    clear_tf(cap.text_frame)
    add_para(cap.text_frame, "系统功能模块图", size=9.2, color=BLUE_DARK, align=PP_ALIGN.CENTER, after=0)

    right = add_panel(slide, Cm(10.0), Cm(2.35), Cm(13.2), Cm(3.5), fill=WHITE, transparency=0.02)
    tf = right.text_frame
    clear_tf(tf)
    add_para(tf, "核心模块", size=14.2, bold=True, color=BLUE_DARK, after=8)
    add_para(tf, "用户管理、运营总览、知识管理、智能问答、预测分析、决策支持。", size=11.3, color=TEXT, after=4)

    bottom_left = add_panel(slide, Cm(10.0), Cm(6.2), Cm(6.2), Cm(4.6), fill=SOFT_BLUE, transparency=0.04)
    tf2 = bottom_left.text_frame
    clear_tf(tf2)
    add_para(tf2, "3 张核心表", size=13.6, bold=True, color=BLUE_DARK, after=8)
    for item in ["sys_user", "kb_document", "ops_daily_metric"]:
        add_para(tf2, item, size=11.5, bullet=True, color=TEXT, after=4)

    bottom_right = add_panel(slide, Cm(17.0), Cm(6.2), Cm(6.2), Cm(4.6), fill=WHITE, transparency=0.02)
    tf3 = bottom_right.text_frame
    clear_tf(tf3)
    add_para(tf3, "4 个关键接口", size=13.6, bold=True, color=BLUE_DARK, after=8)
    for item in ["/api/auth/login", "/api/ops/dashboard", "/api/rag/ask", "/api/ml/predict-*"]:
        add_para(tf3, item, size=10.8, bullet=True, color=TEXT, after=4)


def build_page_show(prs: Presentation) -> None:
    slide = prs.slides.add_slide(prs.slide_layouts[0])
    add_bg(slide, "image4.png")
    add_title_plain(slide, "系统页面展示", "优先展示真实页面，而不是论文大段文字")

    add_pic(slide, "image21.png", Cm(1.0), Cm(2.3), width=Cm(10.8))
    add_pic(slide, "image29.png", Cm(12.0), Cm(2.3), width=Cm(11.0))
    for left, text in [(Cm(1.0), "登录页面"), (Cm(12.0), "运营总览与分析页面")]:
        cap = add_textbox(slide, left, Cm(7.2), Cm(11.0), Cm(0.42))
        clear_tf(cap.text_frame)
        add_para(cap.text_frame, text, size=9.2, color=BLUE_DARK, align=PP_ALIGN.CENTER, after=0)

    note = add_panel(slide, Cm(1.0), Cm(8.05), Cm(22.0), Cm(2.6), fill=WHITE, transparency=0.02)
    clear_tf(note.text_frame)
    add_para(note.text_frame, "页面价值", size=13.8, bold=True, color=BLUE_DARK, after=6)
    add_para(note.text_frame, "系统已完成登录、看板展示和经营状态分析页面，答辩时可直接对应论文实现部分。", size=11.2, color=TEXT, after=0)


def build_qa(prs: Presentation) -> None:
    slide = prs.slides.add_slide(prs.slide_layouts[0])
    add_bg(slide, "image4.png")
    add_title_plain(slide, "知识管理与智能问答", "重点体现 RAG 的知识来源与问答能力")

    pics = [("image25.png", "知识管理"), ("image26.png", "普通问答"), ("image27.png", "深度搜索")]
    xs = [1.0, 8.4, 15.8]
    for (img, label), x in zip(pics, xs):
        add_pic(slide, img, Cm(x), Cm(2.35), width=Cm(6.6))
        cap = add_textbox(slide, Cm(x), Cm(6.65), Cm(6.6), Cm(0.4))
        clear_tf(cap.text_frame)
        add_para(cap.text_frame, label, size=9.0, color=BLUE_DARK, align=PP_ALIGN.CENTER, after=0)

    note = add_panel(slide, Cm(1.0), Cm(7.45), Cm(22.0), Cm(3.2), fill=SOFT_BLUE, transparency=0.05)
    tf = note.text_frame
    clear_tf(tf)
    add_para(tf, "展示要点", size=13.8, bold=True, color=BLUE_DARK, after=6)
    for item in [
        "知识管理页面支持文档新增、查询、更新与删除。",
        "普通问答会返回回答正文与引用信息。",
        "深度搜索面向复杂问题场景，增强回答完整性。",
    ]:
        add_para(tf, item, size=11.0, bullet=True, color=TEXT, after=3)


def build_predict(prs: Presentation) -> None:
    slide = prs.slides.add_slide(prs.slide_layouts[0])
    add_bg(slide, "image4.png")
    add_title_plain(slide, "预测分析与决策支持", "预测结果会继续参与经营判断，而不是只停留在数值层")

    add_pic(slide, "image28.png", Cm(1.0), Cm(2.35), width=Cm(10.0))
    add_pic(slide, "image30.png", Cm(11.7), Cm(2.35), width=Cm(11.3))
    cap1 = add_textbox(slide, Cm(1.0), Cm(6.9), Cm(10.0), Cm(0.4))
    clear_tf(cap1.text_frame)
    add_para(cap1.text_frame, "预测结果展示", size=9.0, color=BLUE_DARK, align=PP_ALIGN.CENTER, after=0)
    cap2 = add_textbox(slide, Cm(11.7), Cm(6.9), Cm(11.3), Cm(0.4))
    clear_tf(cap2.text_frame)
    add_para(cap2.text_frame, "趋势监控页面", size=9.0, color=BLUE_DARK, align=PP_ALIGN.CENTER, after=0)

    note = add_panel(slide, Cm(1.0), Cm(7.6), Cm(22.0), Cm(3.0), fill=WHITE, transparency=0.02)
    tf = note.text_frame
    clear_tf(tf)
    add_para(tf, "页面价值", size=13.8, bold=True, color=BLUE_DARK, after=6)
    add_para(tf, "系统统一展示预测值、趋势变化、风险提示和建议动作，能够和运营总览、知识问答形成闭环。", size=11.1, color=TEXT, after=0)


def build_metrics(prs: Presentation) -> None:
    slide = prs.slides.add_slide(prs.slide_layouts[0])
    add_bg(slide, "image4.png")
    add_title_plain(slide, "模型实验结果总览", "强调“结果可量化”，但不过度堆指标")

    cards = [
        ("入住率预测", ["RMSE 0.0705", "MAE 0.0466", "R² 0.9200"]),
        ("营收预测", ["RMSE 3930.8891", "MAE 2731.3515", "R² 0.9420"]),
        ("订单需求分类", ["Accuracy 0.9099", "Precision 0.8929", "Recall 0.9804", "F1 0.9346"]),
    ]
    xs = [1.0, 8.3, 15.6]
    for (title, items), x in zip(cards, xs):
        panel = add_panel(slide, Cm(x), Cm(2.5), Cm(6.8), Cm(3.8), fill=WHITE, transparency=0.02)
        tf = panel.text_frame
        clear_tf(tf)
        add_para(tf, title, size=13.6, bold=True, color=BLUE_DARK, after=7)
        for item in items:
            add_para(tf, item, size=10.9, bullet=True, color=TEXT, after=3)

    note = add_panel(slide, Cm(1.0), Cm(7.0), Cm(22.0), Cm(3.8), fill=SOFT_BLUE, transparency=0.05)
    tf2 = note.text_frame
    clear_tf(tf2)
    add_para(tf2, "结果解读", size=13.8, bold=True, color=BLUE_DARK, after=6)
    for item in [
        "两项回归任务 R² 均超过 0.92，说明趋势拟合效果较好。",
        "分类任务 Accuracy 超过 0.90，适合承担预警类辅助判断。",
        "目标不是追求生产级极限精度，而是证明系统具备可预测、可联动能力。",
    ]:
        add_para(tf2, item, size=11.0, bullet=True, color=TEXT, after=3)


def build_result_and_test(prs: Presentation) -> None:
    slide = prs.slides.add_slide(prs.slide_layouts[0])
    add_bg(slide, "image4.png")
    add_title_plain(slide, "结果图与系统测试", "一页同时说明实验图与主流程联调结果")

    add_pic(slide, "image31.png", Cm(1.0), Cm(2.3), width=Cm(10.6))
    cap1 = add_textbox(slide, Cm(1.0), Cm(6.75), Cm(10.6), Cm(0.4))
    clear_tf(cap1.text_frame)
    add_para(cap1.text_frame, "入住率预测值与实际值对比", size=8.9, color=BLUE_DARK, align=PP_ALIGN.CENTER, after=0)

    add_pic(slide, "image35.png", Cm(12.0), Cm(2.3), width=Cm(11.0))
    cap2 = add_textbox(slide, Cm(12.0), Cm(6.75), Cm(11.0), Cm(0.4))
    clear_tf(cap2.text_frame)
    add_para(cap2.text_frame, "订单需求分类 ROC 曲线", size=8.9, color=BLUE_DARK, align=PP_ALIGN.CENTER, after=0)

    left = add_panel(slide, Cm(1.0), Cm(7.45), Cm(10.2), Cm(3.2), fill=WHITE, transparency=0.02)
    tf = left.text_frame
    clear_tf(tf)
    add_para(tf, "图表说明", size=13.3, bold=True, color=BLUE_DARK, after=6)
    add_para(tf, "回归图说明趋势拟合能力，ROC 曲线说明分类模型具有较好的区分能力。", size=10.8, color=TEXT, after=0)

    right = add_panel(slide, Cm(12.0), Cm(7.45), Cm(11.0), Cm(3.2), fill=SOFT_BLUE, transparency=0.05)
    tf2 = right.text_frame
    clear_tf(tf2)
    add_para(tf2, "系统测试通过项", size=13.3, bold=True, color=BLUE_DARK, after=6)
    for item in [
        "登录鉴权：成功返回 Token",
        "运营总览：可返回指标与趋势数据",
        "智能问答与三类预测接口均能返回结果",
    ]:
        add_para(tf2, item, size=10.7, bullet=True, color=TEXT, after=3)


def build_summary(prs: Presentation) -> None:
    slide = prs.slides.add_slide(prs.slide_layouts[0])
    add_bg(slide, "image4.png")
    add_title_plain(slide, "总结与展望", "把结论、创新点和后续方向压到同一页")

    a = add_panel(slide, Cm(1.0), Cm(2.35), Cm(7.0), Cm(3.6), fill=WHITE, transparency=0.02)
    clear_tf(a.text_frame)
    add_para(a.text_frame, "总结", size=13.8, bold=True, color=BLUE_DARK, after=8)
    add_para(a.text_frame, "系统已完成核心业务闭环，具备可运行性、可展示性和一定工程实践价值。", size=10.9, color=TEXT, after=0)

    b = add_panel(slide, Cm(8.7), Cm(2.35), Cm(7.0), Cm(3.6), fill=SOFT_BLUE, transparency=0.05)
    clear_tf(b.text_frame)
    add_para(b.text_frame, "创新点", size=13.8, bold=True, color=BLUE_DARK, after=8)
    for item in [
        "整合看板、知识、问答、预测与决策支持。",
        "强调问答与预测结果的可解释性。",
    ]:
        add_para(b.text_frame, item, size=10.8, bullet=True, color=TEXT, after=3)

    c = add_panel(slide, Cm(16.4), Cm(2.35), Cm(7.0), Cm(3.6), fill=WHITE, transparency=0.02)
    clear_tf(c.text_frame)
    add_para(c.text_frame, "展望", size=13.8, bold=True, color=BLUE_DARK, after=8)
    for item in [
        "接入真实 PMS 数据。",
        "增强知识在线更新与风险分析能力。",
    ]:
        add_para(c.text_frame, item, size=10.8, bullet=True, color=TEXT, after=3)

    d = add_panel(slide, Cm(1.0), Cm(6.7), Cm(22.3), Cm(4.0), fill=BLUE, transparency=0.02)
    tf = d.text_frame
    clear_tf(tf)
    add_para(tf, "答辩表达建议", size=14.0, bold=True, color=WHITE, after=8)
    add_para(tf, "当前版本重点证明“链路已经打通、页面已经实现、结果可以验证”，未来优化重点放在真实数据接入与长期迭代。", size=11.1, color=WHITE, after=0)


def build_thanks(prs: Presentation) -> None:
    slide = prs.slides.add_slide(prs.slide_layouts[0])
    add_bg(slide, "image5.png")
    tb = add_textbox(slide, Cm(5.0), Cm(4.1), Cm(15.0), Cm(2.4))
    clear_tf(tb.text_frame)
    add_para(tb.text_frame, "感谢各位老师聆听", size=24, bold=True, color=BLUE_DARK, align=PP_ALIGN.CENTER, after=6)
    add_para(tb.text_frame, "恳请批评指正", size=13.5, color=MUTED, align=PP_ALIGN.CENTER, after=0)
    info = add_textbox(slide, Cm(6.1), Cm(8.6), Cm(13.0), Cm(0.9))
    clear_tf(info.text_frame)
    add_para(info.text_frame, "刘洋 | 人工智能 3 班 | 指导教师：刘会芝", size=10.8, color=MUTED, align=PP_ALIGN.CENTER, after=4)
    add_para(info.text_frame, "基于SpringBoot与RAG的酒店运营支持系统设计与实现", size=11.8, bold=True, color=BLUE_DARK, align=PP_ALIGN.CENTER, after=0)


def build() -> None:
    prs = Presentation(str(TEMPLATE))
    delete_all_slides(prs)

    build_cover(prs)
    build_toc(prs)
    build_section(prs, "研究背景与系统设计", "Problem, Goal and Architecture")
    build_background_gap(prs)
    build_goal_arch(prs)
    build_module_data(prs)
    build_section(prs, "系统实现与页面展示", "Implementation and Demonstration")
    build_page_show(prs)
    build_qa(prs)
    build_predict(prs)
    build_section(prs, "实验结果与测试验证", "Experiment and Validation")
    build_metrics(prs)
    build_result_and_test(prs)
    build_summary(prs)
    build_thanks(prs)

    validate(prs)
    prs.save(str(OUTPUT))
    print(OUTPUT)


if __name__ == "__main__":
    build()
